"""
Generate the Autonomous Social Media Content Studio demo presentation.
Run with:  uv run python generate_ppt.py
Output  :  demo_presentation.pptx  (in the project root)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)   # electric cyan
ACCENT2    = RGBColor(0x7B, 0x2F, 0xFF)   # purple
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD6, 0xE0)
GREEN      = RGBColor(0x00, 0xE5, 0x96)
YELLOW     = RGBColor(0xFF, 0xD1, 0x66)
ORANGE     = RGBColor(0xFF, 0x6B, 0x35)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, title=None, title_size=18,
                   title_color=ACCENT, bullet_color=WHITE,
                   bullet_char="▸ ", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = bullet_color
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def section_header(slide, label, color=ACCENT):
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), color)
    add_text(slide, label, Inches(0.35), Inches(0.12), Inches(12), Inches(0.45),
             font_size=11, color=color, bold=True)


def slide_title(slide, title, subtitle=None):
    add_text(slide, title, Inches(0.35), Inches(0.6), Inches(12.5), Inches(0.8),
             font_size=34, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(1.35), Inches(12.5), Inches(0.4),
                 font_size=17, color=LIGHT_GRAY)


def fill_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════════
def slide_cover():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)

    # gradient-like left panel
    add_rect(slide, Inches(0), Inches(0), Inches(4.8), SLIDE_H, ACCENT2)
    add_rect(slide, Inches(4.6), Inches(0), Inches(0.4), SLIDE_H, RGBColor(0x3A, 0x1A, 0x7A))

    # emoji / icon area
    add_text(slide, "🤖", Inches(0.5), Inches(0.6), Inches(4), Inches(1.2),
             font_size=60, align=PP_ALIGN.CENTER)

    add_text(slide, "CAPSTONE PROJECT", Inches(5.2), Inches(0.7), Inches(7.8), Inches(0.5),
             font_size=12, color=ACCENT, bold=True, italic=True)

    add_text(slide,
             "Autonomous Social Media\nContent Studio",
             Inches(5.2), Inches(1.2), Inches(7.8), Inches(2.0),
             font_size=36, bold=True, color=WHITE)

    add_text(slide,
             "A Multi-Agent AI System that plans, writes, reviews\n"
             "and refines platform-native social media content — autonomously.",
             Inches(5.2), Inches(3.1), Inches(7.8), Inches(1.2),
             font_size=16, color=LIGHT_GRAY)

    # divider
    add_rect(slide, Inches(5.2), Inches(4.4), Inches(7.0), Inches(0.04), ACCENT)

    add_text(slide, "Contributors", Inches(5.2), Inches(4.55), Inches(7.0), Inches(0.4),
             font_size=13, color=ACCENT, bold=True)
    add_text(slide, "Himanshu Jain  •  Prashanchal Sharma",
             Inches(5.2), Inches(4.9), Inches(7.0), Inches(0.4),
             font_size=15, color=WHITE, bold=True)

    add_text(slide, "Technology Stack: LangGraph  •  LangChain  •  Azure OpenAI  •  FastAPI",
             Inches(5.2), Inches(5.5), Inches(7.8), Inches(0.4),
             font_size=12, color=LIGHT_GRAY)
    add_text(slide, "April 2026",
             Inches(5.2), Inches(6.8), Inches(3), Inches(0.4),
             font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Solution
# ═══════════════════════════════════════════════════════════════════════════════
def slide_problem_solution():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "PROBLEM & SOLUTION")
    slide_title(slide, "The Problem We Solve")

    # Problem box
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.8), Inches(4.8),
             RGBColor(0x1A, 0x0A, 0x2E))
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.8), Inches(0.5), ORANGE)
    add_text(slide, "😩  THE PROBLEM", Inches(0.5), Inches(2.05), Inches(5.6), Inches(0.42),
             font_size=15, bold=True, color=WHITE)
    problems = [
        "Creating quality content for 3 platforms daily is exhausting",
        "Each platform needs a different tone, format and character limit",
        "Manual research + writing + reviewing takes hours per campaign",
        "Inconsistent quality when working under time pressure",
        "Hard to adapt a single idea across LinkedIn, X and Instagram",
    ]
    add_bullet_box(slide, problems,
                   Inches(0.5), Inches(2.6), Inches(5.5), Inches(4.0),
                   font_size=14, bullet_color=LIGHT_GRAY, bullet_char="✗  ")

    # Solution box
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(6.1), Inches(4.8),
             RGBColor(0x02, 0x1A, 0x1A))
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(6.1), Inches(0.5), GREEN)
    add_text(slide, "🚀  OUR SOLUTION", Inches(6.9), Inches(2.05), Inches(5.9), Inches(0.42),
             font_size=15, bold=True, color=WHITE)
    solutions = [
        "Autonomous multi-agent system powered by LangGraph",
        "Planner agent researches trends & builds a content strategy",
        "3 specialist agents write platform-native posts in parallel",
        "Reviewer agent scores quality & flags issues automatically",
        "Reflection loop refines only failing posts — not all platforms",
    ]
    add_bullet_box(slide, solutions,
                   Inches(6.9), Inches(2.6), Inches(5.8), Inches(4.0),
                   font_size=14, bullet_color=LIGHT_GRAY, bullet_char="✓  ")

    # Arrow in the middle
    add_text(slide, "→", Inches(6.1), Inches(4.1), Inches(0.65), Inches(0.6),
             font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Strategy & Multi-Agent Choice
# ═══════════════════════════════════════════════════════════════════════════════
def slide_strategy():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "PROJECT STRATEGY")
    slide_title(slide, "Strategy & Multi-Agent Pattern",
                "Why Plan-Act-Check-Reflect with LangGraph?")

    # Left: strategy
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.9), Inches(5.1),
             RGBColor(0x10, 0x20, 0x38))
    add_text(slide, "🎯  Strategy", Inches(0.5), Inches(2.1), Inches(5.7), Inches(0.45),
             font_size=16, bold=True, color=ACCENT)
    strat = [
        "Input: topic text, uploaded file, or a blog URL",
        "DuckDuckGo search enriches planning with live data",
        "Single planner creates one coherent strategy",
        "3 agents write simultaneously (LinkedIn / X / Instagram)",
        "LLM reviewer scores each post 0-10 per quality rubric",
        "Only failing platforms (<7.0) are retried in the next loop",
        "Loop exits when score ≥ 8.8, no improvement, or max 3 iters",
    ]
    add_bullet_box(slide, strat, Inches(0.5), Inches(2.6), Inches(5.7), Inches(4.3),
                   font_size=13, bullet_color=WHITE)

    # Right: why this pattern
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.1),
             RGBColor(0x10, 0x20, 0x38))
    add_text(slide, "🧠  Why Plan-Act-Check-Reflect?", Inches(6.9), Inches(2.1), Inches(5.9), Inches(0.45),
             font_size=16, bold=True, color=ACCENT2)
    why = [
        "Plan-Act-Check gives structure to an open-ended task",
        "Parallel fan-out cuts latency vs sequential generation",
        "Self-correction loop removes need for human review rounds",
        "Targeted reflect avoids re-running already-passing agents",
        "LangGraph state machine handles fan-in/fan-out cleanly",
        "Stateful graph makes iteration count & scores observable",
        "Azure OpenAI prompt caching cuts cost on repeated calls",
    ]
    add_bullet_box(slide, why, Inches(6.9), Inches(2.6), Inches(5.8), Inches(4.3),
                   font_size=13, bullet_color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture & Libraries
# ═══════════════════════════════════════════════════════════════════════════════
def slide_architecture():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "ARCHITECTURE")
    slide_title(slide, "Project Architecture & Tech Stack")

    # Architecture diagram (text-based)
    arch_lines = [
        "┌──────────────────────────────────────────────────────────────────────┐",
        "│                        FastAPI  (REST API)                           │",
        "│   POST /upload  ─────────────────────┐   POST /generate ──────────┐ │",
        "│   Saves file, returns path            │   Starts LangGraph workflow│ │",
        "└──────────────────────────────────────────────────────────────────────┘",
        "                               │",
        "              ┌────────────────▼────────────────┐",
        "              │       LangGraph State Machine    │",
        "              │  START → PLAN → [fan-out]        │",
        "              │  [LinkedIn│X│Instagram] → REVIEW │",
        "              │  REVIEW → approve? → FORMAT→END  │",
        "              │          ↘ reflect? → REFLECT    │",
        "              └─────────────────────────────────-┘",
    ]
    add_text(slide, "\n".join(arch_lines),
             Inches(0.35), Inches(1.85), Inches(8.5), Inches(3.5),
             font_size=9.5, color=ACCENT)

    # Libs panel
    libs = [
        ("langgraph ≥ 0.2",   "State machine & graph orchestration",  ACCENT),
        ("langchain ≥ 0.3",   "LLM abstractions & tool protocol",     ACCENT2),
        ("langchain-openai",   "Azure OpenAI chat model adapter",      GREEN),
        ("fastapi ≥ 0.115",   "REST API server (upload + generate)",  YELLOW),
        ("uvicorn",            "ASGI server for FastAPI",              LIGHT_GRAY),
        ("pydantic ≥ 2.0",    "Request/response models & validation", LIGHT_GRAY),
        ("duckduckgo-search",  "Live web research for planner agent",  ORANGE),
        ("pypdf ≥ 4.0",       "PDF text extraction for file input",   ORANGE),
        ("python-multipart",   "Multipart file upload support",        LIGHT_GRAY),
        ("python-dotenv",      "Environment variable management",      LIGHT_GRAY),
    ]
    add_rect(slide, Inches(8.9), Inches(1.85), Inches(4.0), Inches(5.35),
             RGBColor(0x10, 0x20, 0x38))
    add_text(slide, "📦  Dependencies", Inches(9.0), Inches(1.9), Inches(3.8), Inches(0.4),
             font_size=13, bold=True, color=ACCENT)
    y = Inches(2.35)
    for pkg, desc, col in libs:
        add_text(slide, pkg, Inches(9.05), y, Inches(2.0), Inches(0.32),
                 font_size=10, bold=True, color=col)
        add_text(slide, desc, Inches(9.05), y + Inches(0.22), Inches(3.7), Inches(0.28),
                 font_size=9, color=LIGHT_GRAY)
        y += Inches(0.49)

    # Folder structure
    add_text(slide, "📁  app/\n  core/  (llm, logger)\n  nodes/ (plan, act ×3, check, reflect, format)\n  graph/ (graph, routing)\n  prompts/ (per-agent system + human messages)\n  tools/  (web_search, read_file, read_url,\n           char_counter, guidelines_checker)\n  state/  (ContentState TypedDict)\n  models/ (Pydantic request/response)\n  main.py (FastAPI app)",
             Inches(0.35), Inches(5.35), Inches(8.4), Inches(1.9),
             font_size=10, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LangGraph Workflow Diagram
# ═══════════════════════════════════════════════════════════════════════════════
def slide_workflow():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "LANGGRAPH WORKFLOW")
    slide_title(slide, "Visual Workflow — Plan · Act · Check · Reflect Loop")

    LW = Pt(2.5)  # line weight

    # ── line / arrow helpers ──────────────────────────────────────────────────
    def hline(x1, y, x2, color=LIGHT_GRAY):
        w = Inches(x2 - x1)
        if w <= 0:
            return
        add_rect(slide, Inches(x1), Inches(y) - LW // 2, w, LW, color)

    def vline(x, y1, y2, color=LIGHT_GRAY):
        h = Inches(y2 - y1)
        if h <= 0:
            return
        add_rect(slide, Inches(x) - LW // 2, Inches(y1), LW, h, color)

    def arrowD(x, y, color=LIGHT_GRAY):
        add_text(slide, "▼", Inches(x - .15), Inches(y), Inches(.3), Inches(.22),
                 font_size=8, color=color, align=PP_ALIGN.CENTER)

    def arrowR(x, y, color=LIGHT_GRAY):
        add_text(slide, "▶", Inches(x), Inches(y - .11), Inches(.22), Inches(.22),
                 font_size=8, color=color, align=PP_ALIGN.CENTER)

    def fnode(label, cx, cy, w, h, bg, fs=12):
        add_rect(slide, Inches(cx - w / 2), Inches(cy - h / 2), Inches(w), Inches(h), bg)
        add_text(slide, label, Inches(cx - w / 2), Inches(cy - h / 2), Inches(w), Inches(h),
                 font_size=fs, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    def label(txt, x, y, w=2.5, color=LIGHT_GRAY, fs=9, bold=False, italic=False):
        add_text(slide, txt, Inches(x), Inches(y), Inches(w), Inches(.28),
                 font_size=fs, color=color, bold=bold, italic=italic)

    # ── layout (all coords in inches, nodes positioned by CENTER) ─────────────
    CX, NW, NH = 6.17, 2.6, 0.50          # main column  center-x / width / height
    AW = 2.0                                # agent node width

    Y_SRT = 1.90   # START
    Y_PLN = 2.62   # PLAN
    Y_FO  = 3.14   # fan-out rail
    Y_AGT = 3.66   # platform agents
    Y_FI  = 4.18   # fan-in rail
    Y_CHK = 4.72   # CHECK
    Y_DEC = 5.40   # decision
    Y_SID = 6.12   # REFLECT / FORMAT
    Y_END = 6.88   # END

    A_LI, A_X, A_IG = 2.4, CX, 9.9        # agent center-x
    REF_X = 1.9                             # REFLECT center-x
    FMT_X = 11.0                            # FORMAT / END center-x
    LPX   = 0.32                            # loop gutter x

    # ── draw lines FIRST so nodes appear on top ───────────────────────────────

    # START → PLAN
    vline(CX, Y_SRT + NH / 2, Y_PLN - NH / 2)
    arrowD(CX, Y_PLN - NH / 2 - .20)

    # PLAN → fan-out rail
    vline(CX, Y_PLN + NH / 2, Y_FO)

    # fan-out horizontal rail (A_LI to A_IG)
    hline(A_LI, Y_FO, A_IG)

    # fan-out rail → each agent
    for ax in [A_LI, A_X, A_IG]:
        vline(ax, Y_FO, Y_AGT - NH / 2)
        arrowD(ax, Y_AGT - NH / 2 - .20)

    # each agent → fan-in rail
    for ax in [A_LI, A_X, A_IG]:
        vline(ax, Y_AGT + NH / 2, Y_FI)

    # fan-in horizontal rail
    hline(A_LI, Y_FI, A_IG)

    # fan-in rail → CHECK
    vline(CX, Y_FI, Y_CHK - NH / 2)
    arrowD(CX, Y_CHK - NH / 2 - .20)

    # CHECK → decision
    vline(CX, Y_CHK + NH / 2, Y_DEC - NH / 2)
    arrowD(CX, Y_DEC - NH / 2 - .20)

    # YES path: right corner → FORMAT → END  (green)
    hline(CX + NW / 2, Y_DEC, FMT_X, GREEN)
    vline(FMT_X, Y_DEC, Y_SID - NH / 2, GREEN)
    arrowD(FMT_X, Y_SID - NH / 2 - .20, GREEN)
    vline(FMT_X, Y_SID + NH / 2, Y_END - NH / 2, GREEN)
    arrowD(FMT_X, Y_END - NH / 2 - .20, GREEN)

    # NO path: left corner → REFLECT  (orange)
    hline(REF_X, Y_DEC, CX - NW / 2, ORANGE)
    vline(REF_X, Y_DEC, Y_SID - NH / 2, ORANGE)
    arrowD(REF_X, Y_SID - NH / 2 - .20, ORANGE)

    # Loop: REFLECT bottom → left gutter → up → PLAN left  (orange)
    hline(LPX, Y_SID, REF_X - AW / 2, ORANGE)
    vline(LPX, Y_PLN, Y_SID, ORANGE)
    hline(LPX, Y_PLN, CX - NW / 2, ORANGE)
    arrowR(CX - NW / 2 - .22, Y_PLN, ORANGE)

    # ── nodes (rendered on top of lines) ─────────────────────────────────────
    fnode("START",        CX,    Y_SRT, 1.6,  NH * .9, ACCENT2,                  11)
    fnode("🗺  PLAN",     CX,    Y_PLN, NW,   NH,      RGBColor(0x1A, 0x6B, 0xBF))
    fnode("LinkedIn",     A_LI,  Y_AGT, AW,   NH,      RGBColor(0x00, 0x77, 0xB5))
    fnode("X / Twitter",  A_X,   Y_AGT, AW,   NH,      RGBColor(0x22, 0x22, 0x22))
    fnode("Instagram",    A_IG,  Y_AGT, AW,   NH,      RGBColor(0xE1, 0x30, 0x6C))
    fnode("🔍  CHECK",   CX,    Y_CHK, NW,   NH,      RGBColor(0x8B, 0x45, 0x13))
    fnode("score ≥ 8.8?", CX,    Y_DEC, NW,   NH,      RGBColor(0x2A, 0x2A, 0x50))
    fnode("🔄  REFLECT",  REF_X, Y_SID, AW,   NH,      ACCENT2)
    fnode("✅  FORMAT",   FMT_X, Y_SID, AW,   NH,      GREEN)
    fnode("END",          FMT_X, Y_END, 1.6,  NH * .9, ACCENT2,                  11)

    # ── line labels ───────────────────────────────────────────────────────────
    label("fan-out (parallel)", CX + NW / 2 + .12, Y_FO - .27, color=ACCENT, italic=True)
    label("fan-in",             CX + NW / 2 + .12, Y_FI - .27, color=ACCENT, italic=True)
    label("YES ✓",   CX + NW / 2 + .12, Y_DEC - .30, color=GREEN,  fs=10, bold=True)
    label("NO  ✗",   REF_X + AW / 2 + .08, Y_DEC - .30, color=ORANGE, fs=10, bold=True)
    label("retry",   LPX + .06, (Y_PLN + Y_SID) / 2 - .14, w=.8, color=ORANGE, italic=True)

    # ── routing rules panel (top-right, away from diagram) ───────────────────
    rules = [
        "Quality pass  : overall score ≥ 8.8",
        "Platform retry: per-platform score < 7.0",
        "Max iterations: 3  (configurable)",
        "Stagnation exit: no improvement → approve",
    ]
    add_rect(slide, Inches(9.5), Inches(1.75), Inches(3.6), Inches(1.5),
             RGBColor(0x10, 0x20, 0x38))
    add_bullet_box(slide, rules, Inches(9.6), Inches(1.75), Inches(3.45), Inches(1.5),
                   font_size=9.5, title="Routing Rules", title_size=10.5,
                   title_color=YELLOW, bullet_color=LIGHT_GRAY, bullet_char="  ")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ContentState
# ═══════════════════════════════════════════════════════════════════════════════
def slide_state():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "LANGGRAPH STATE")
    slide_title(slide, "ContentState — Shared Agent Memory")

    groups = [
        ("📥  Input",
         ["input_content  (str)  — topic/brief text",
          "source_url     (str)  — blog post URL (optional)",
          "uploaded_file_path (str) — PDF/TXT file (optional)"],
         ACCENT, Inches(0.4), Inches(2.0)),
        ("🗺  Plan",
         ["content_plan  (str)  — strategy produced by Planner",
          "search_results (str) — DuckDuckGo research context"],
         ACCENT2, Inches(0.4), Inches(3.45)),
        ("✍️  Act",
         ["posts  (Dict[str,str], merge_dicts reducer)",
          "  → {linkedin, x, instagram} written in parallel"],
         GREEN, Inches(0.4), Inches(4.6)),
        ("🔍  Check",
         ["feedback       (Dict, merge_dicts)  — per-platform review",
          "overall_score  (float)  — aggregated quality score",
          "previous_score (float)  — prior iteration score"],
         YELLOW, Inches(6.8), Inches(2.0)),
        ("🔄  Reflect",
         ["platforms_to_retry  (list) — failed platforms only"],
         ORANGE, Inches(6.8), Inches(3.35)),
        ("✅  Format",
         ["final_output   (str)  — publish-ready JSON package",
          "approval_status (bool)"],
         RGBColor(0x00, 0xAA, 0x66), Inches(6.8), Inches(4.3)),
        ("⚙️  Control",
         ["iteration_count (int)",
          "max_iterations  (int, default 3)",
          "history         (list, append_list reducer)",
          "token_usage     (dict, sum_usage reducer) — accumulated tokens"],
         LIGHT_GRAY, Inches(6.8), Inches(5.2)),
    ]

    for title, items, color, lft, tp in groups:
        box_h = Inches(0.45 + len(items) * 0.38)
        add_rect(slide, lft, tp, Inches(5.9), box_h, RGBColor(0x10, 0x20, 0x38))
        add_rect(slide, lft, tp, Inches(5.9), Inches(0.08), color)
        add_text(slide, title, lft + Inches(0.1), tp + Inches(0.1), Inches(5.7), Inches(0.4),
                 font_size=13, bold=True, color=color)
        for i, item in enumerate(items):
            add_text(slide, "  " + item, lft + Inches(0.15),
                     tp + Inches(0.5 + i * 0.38), Inches(5.6), Inches(0.36),
                     font_size=11, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Planner Agent
# ═══════════════════════════════════════════════════════════════════════════════
def slide_planner():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "AGENTS — PLANNER")
    slide_title(slide, "Planner Agent  🗺", "Builds the content strategy before any post is written")

    add_rect(slide, Inches(0.4), Inches(1.95), Inches(5.9), Inches(5.2),
             RGBColor(0x0A, 0x20, 0x40))

    add_text(slide, "What it does", Inches(0.5), Inches(2.0), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=ACCENT)
    planner_what = [
        "Reads input (text / uploaded file / URL)",
        "Calls web_search ×2 — topic research + trending content",
        "Runs LLM call with research to build 9-section content plan",
        "Defines audience, core message, hashtag themes per platform",
        "Also handles REPLAN after the Check→Reflect loop",
    ]
    add_bullet_box(slide, planner_what, Inches(0.5), Inches(2.45), Inches(5.6), Inches(2.2),
                   font_size=13, bullet_color=WHITE)

    add_text(slide, "Plan sections", Inches(0.5), Inches(4.75), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=ACCENT)
    sections = [
        "1 Target audience   5 LinkedIn strategy",
        "2 Core message      6 X (Twitter) strategy",
        "3 Key facts/data    7 Instagram strategy",
        "4 Trending angles   8 Key phrases   9 Avoid list",
    ]
    add_bullet_box(slide, sections, Inches(0.5), Inches(5.2), Inches(5.6), Inches(1.8),
                   font_size=12, bullet_color=LIGHT_GRAY, bullet_char="   ")

    # Right: source content flow
    add_rect(slide, Inches(6.8), Inches(1.95), Inches(6.1), Inches(5.2),
             RGBColor(0x0A, 0x20, 0x40))
    add_text(slide, "Source content flow", Inches(6.9), Inches(2.0), Inches(5.9), Inches(0.4),
             font_size=14, bold=True, color=ACCENT2)

    flows = [
        ("📝  text input", "Direct brief from user", WHITE),
        ("📎  uploaded file", "PDF, TXT, MD, CSV via /upload endpoint\nread_uploaded_file tool extracts text", WHITE),
        ("🌐  URL input", "read_url tool fetches & strips HTML\nfrom any blog post or article", WHITE),
    ]
    y = Inches(2.5)
    for icon, desc, col in flows:
        add_rect(slide, Inches(6.9), y, Inches(5.7), Inches(0.85),
                 RGBColor(0x15, 0x2A, 0x50))
        add_text(slide, icon, Inches(7.0), y + Inches(0.05), Inches(2.5), Inches(0.38),
                 font_size=13, bold=True, color=ACCENT)
        add_text(slide, desc, Inches(7.0), y + Inches(0.42), Inches(5.5), Inches(0.4),
                 font_size=11, color=LIGHT_GRAY)
        y += Inches(1.0)

    add_text(slide, "⚡  Prompt caching: static PLAN_SYSTEM\n    prepended with 1 100-token agency brief\n    → Azure cache HIT from 2nd request onward",
             Inches(6.9), Inches(5.6), Inches(5.8), Inches(0.9),
             font_size=11, color=YELLOW)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Act Agents
# ═══════════════════════════════════════════════════════════════════════════════
def slide_act():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "AGENTS — ACT")
    slide_title(slide, "Act Agents — Platform Writers  ✍️",
                "Three specialist agents write posts in parallel (LangGraph fan-out)")

    platforms = [
        ("LinkedIn Agent", RGBColor(0x00, 0x77, 0xB5),
         ["Professional, thought-leadership tone",
          "900–1 400 character body recommended",
          "Hook in first 2 lines (≤ 210 chars)",
          "3–5 hashtags appended at the end",
          "0–2 emoji; no URLs in body text",
          "Structure: Hook → Story/Data → Insight → CTA"],
         "3 000 char limit"),
        ("X / Twitter Agent", RGBColor(0x15, 0x1F, 0x2B),
         ["Punchy hook — first 3 words must earn attention",
          "HARD 280-char limit (enforced by character_counter tool)",
          "Max 2 hashtags (count toward 280 chars)",
          "Links count as 23 chars (t.co shortener)",
          "1–2 emoji maximum; text-forward style",
          "character_counter tool trims at word boundary"],
         "280 char HARD limit"),
        ("Instagram Agent", RGBColor(0xC1, 0x32, 0x84),
         ["Aspirational, community-driven 'you/your' voice",
          "First 125 chars appear before the 'more' fold",
          "5–15 hashtags after 5 blank lines",
          "2–6 emoji for warmth and line breaks",
          "Must end with a question or CTA",
          "Structure: Hook → Story → Value → CTA → Tags"],
         "2 200 char limit"),
    ]

    for i, (name, color, points, limit) in enumerate(platforms):
        x = Inches(0.3 + i * 4.35)
        add_rect(slide, x, Inches(1.9), Inches(4.1), Inches(5.3), RGBColor(0x10, 0x20, 0x38))
        add_rect(slide, x, Inches(1.9), Inches(4.1), Inches(0.55), color)
        add_text(slide, name, x + Inches(0.1), Inches(1.93), Inches(3.9), Inches(0.45),
                 font_size=15, bold=True, color=WHITE)
        add_text(slide, limit, x + Inches(0.1), Inches(2.5), Inches(3.9), Inches(0.3),
                 font_size=10, color=YELLOW, italic=True)
        add_bullet_box(slide, points, x + Inches(0.1), Inches(2.85), Inches(3.85), Inches(4.0),
                       font_size=11, bullet_color=LIGHT_GRAY)

    # Fan-out note
    add_rect(slide, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.35),
             RGBColor(0x10, 0x20, 0x38))
    add_text(slide,
             "All 3 agents run in PARALLEL via LangGraph fan-out  •  "
             "Selective retry: only failed platforms re-run after Reflect  •  "
             "Feedback from previous iteration injected into each agent's prompt",
             Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
             font_size=10, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Check Agent
# ═══════════════════════════════════════════════════════════════════════════════
def slide_check():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "AGENTS — CHECK")
    slide_title(slide, "Check Agent — Quality Reviewer  🔍",
                "Scores every post against the content plan and quality rubric")

    # Left: what it does
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.9), Inches(5.2),
             RGBColor(0x0A, 0x20, 0x40))
    add_text(slide, "Scoring criteria (per platform)", Inches(0.5), Inches(2.05), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=ACCENT)
    criteria = [
        "Alignment with the content plan",
        "Tone & format appropriateness for platform",
        "Clarity and engagement quality",
        "Grammar and polish",
    ]
    add_bullet_box(slide, criteria, Inches(0.5), Inches(2.5), Inches(5.6), Inches(1.8),
                   font_size=13, bullet_color=WHITE)

    add_text(slide, "Also runs deterministic tool check", Inches(0.5), Inches(4.4), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=ACCENT)
    tool_check = [
        "content_guidelines_checker tool runs on all 3 posts",
        "Detects brand safety violations (hate speech, explicit content)",
        "Detects over-promotional language",
        "Violations appended to feedback — agent must fix next round",
    ]
    add_bullet_box(slide, tool_check, Inches(0.5), Inches(4.85), Inches(5.6), Inches(2.0),
                   font_size=12, bullet_color=LIGHT_GRAY)

    # Right: JSON output
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.2),
             RGBColor(0x06, 0x18, 0x10))
    add_text(slide, "Output JSON (structured)", Inches(6.9), Inches(2.05), Inches(5.9), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
    json_ex = '''{
  "linkedin":  {
    "score": 8,
    "issues": "Hook could be stronger",
    "suggestions": "Open with a statistic"
  },
  "x":         { "score": 9,  "issues": "None",
                 "suggestions": "Good as-is" },
  "instagram": { "score": 6,
    "issues": "Missing CTA",
    "suggestions": "Add a question at the end"
  },
  "overall_score": 7.67
}'''
    add_text(slide, json_ex, Inches(6.9), Inches(2.5), Inches(5.8), Inches(3.0),
             font_size=10.5, color=GREEN)

    add_text(slide, "Routing decision after check:", Inches(6.9), Inches(5.55), Inches(5.8), Inches(0.35),
             font_size=12, bold=True, color=YELLOW)
    routing = [
        "overall ≥ 8.8         → APPROVE → Format",
        "score stagnated       → APPROVE (avoid wasted loops)",
        "max_iterations hit    → APPROVE (force exit)",
        "otherwise             → REFLECT (targeted replan)",
    ]
    add_bullet_box(slide, routing, Inches(6.9), Inches(5.95), Inches(5.8), Inches(1.2),
                   font_size=11, bullet_color=WHITE, bullet_char="  ")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Reflect Agent
# ═══════════════════════════════════════════════════════════════════════════════
def slide_reflect():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "AGENTS — REFLECT")
    slide_title(slide, "Reflect Agent — Targeted Replanner  🔄",
                "Only retries platforms that failed — not the entire workflow")

    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.9), Inches(5.2),
             RGBColor(0x1A, 0x08, 0x30))
    add_text(slide, "How it works", Inches(0.5), Inches(2.05), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=ACCENT2)
    how = [
        "Reads per-platform scores from Check feedback",
        "Identifies failing platforms (score < 7.0 threshold)",
        "If none below threshold, retries the single lowest scorer",
        "Calls LLM with failing feedback → produces revised plan",
        "Sets platforms_to_retry in state",
        "Fan-out in next iteration only dispatches failing agents",
    ]
    add_bullet_box(slide, how, Inches(0.5), Inches(2.5), Inches(5.6), Inches(2.8),
                   font_size=13, bullet_color=WHITE)

    add_text(slide, "Why targeted, not full replan?", Inches(0.5), Inches(5.4), Inches(5.7), Inches(0.4),
             font_size=13, bold=True, color=ACCENT2)
    why = [
        "Saves 2 LLM calls per iteration (passing platforms skip)",
        "Avoids regressing already-good posts",
        "Prevents oscillation (platforms swapping pass/fail)",
    ]
    add_bullet_box(slide, why, Inches(0.5), Inches(5.85), Inches(5.6), Inches(1.3),
                   font_size=12, bullet_color=LIGHT_GRAY)

    # Right: before / after
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.2),
             RGBColor(0x1A, 0x08, 0x30))
    add_text(slide, "Iteration example", Inches(6.9), Inches(2.05), Inches(5.9), Inches(0.4),
             font_size=14, bold=True, color=ACCENT2)

    example = [
        ("Iter 1 — CHECK output:",
         "LinkedIn: 8.5   X: 9.0   Instagram: 5.5   overall: 7.67", LIGHT_GRAY),
        ("REFLECT decision:",
         "Only Instagram failed (5.5 < 7.0 threshold)", ORANGE),
        ("platforms_to_retry = ['instagram']",
         "LinkedIn & X agents are SKIPPED", YELLOW),
        ("Iter 2 — CHECK output:",
         "Instagram: 8.0   overall: 8.5 → APPROVED ✓", GREEN),
    ]
    y = Inches(2.5)
    for label, detail, col in example:
        add_text(slide, label, Inches(6.9), y, Inches(5.8), Inches(0.3),
                 font_size=12, bold=True, color=col)
        add_text(slide, detail, Inches(6.9), y + Inches(0.3), Inches(5.8), Inches(0.35),
                 font_size=11, color=WHITE)
        y += Inches(0.8)

    add_rect(slide, Inches(6.8), Inches(6.1), Inches(6.1), Inches(0.9),
             RGBColor(0x10, 0x25, 0x10))
    add_text(slide,
             "Result: 5 LLM calls instead of 7 in 2 iterations\n"
             "(saved 2 × Act calls for passing platforms)",
             Inches(6.9), Inches(6.15), Inches(5.9), Inches(0.75),
             font_size=12, color=GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Tools
# ═══════════════════════════════════════════════════════════════════════════════
def slide_tools():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "TOOLS")
    slide_title(slide, "Tools — Agent Capabilities  🔧",
                "LangChain @tool functions — deterministic helpers for agent nodes")

    tools_data = [
        ("🌐  web_search",
         "DuckDuckGo live search",
         ["Used by Planner agent (×2 per generation)",
          "Fetches topic research + trending content",
          "Returns title, URL, snippet for each result",
          "Caps at 10 results to control payload size"],
         ACCENT, Inches(0.4), Inches(2.0)),
        ("📎  read_uploaded_file",
         "Extract text from uploaded file",
         ["Supports .txt, .md, .rst, .csv, .pdf",
          "PDF extraction via pypdf library",
          "Truncates to 8 000 chars to fit prompt budget",
          "Error-safe: guards against malformed files"],
         ACCENT2, Inches(0.4), Inches(4.5)),
        ("🌐  read_url",
         "Fetch & strip HTML from a URL",
         ["Uses stdlib urllib — no heavy dependencies",
          "Strips <script>, <nav>, <footer> tags",
          "Returns clean article body text",
          "Truncates to 8 000 chars"],
         GREEN, Inches(6.8), Inches(2.0)),
        ("📏  character_counter",
         "Enforce X/Twitter 280-char limit",
         ["Deterministic — no LLM needed",
          "Trims at word boundary (not mid-word)",
          "Returns trimmed flag + final text",
          "Used by X agent after every generation"],
         YELLOW, Inches(6.8), Inches(4.5)),
    ]

    for name, subtitle, points, color, lft, tp in tools_data:
        add_rect(slide, lft, tp, Inches(5.9), Inches(2.25), RGBColor(0x10, 0x20, 0x38))
        add_rect(slide, lft, tp, Inches(5.9), Inches(0.08), color)
        add_text(slide, name, lft + Inches(0.1), tp + Inches(0.1), Inches(5.6), Inches(0.38),
                 font_size=14, bold=True, color=color)
        add_text(slide, subtitle, lft + Inches(0.1), tp + Inches(0.45), Inches(5.6), Inches(0.28),
                 font_size=11, color=LIGHT_GRAY, italic=True)
        add_bullet_box(slide, points, lft + Inches(0.1), tp + Inches(0.75), Inches(5.6), Inches(1.4),
                       font_size=11, bullet_color=WHITE)

    # 5th tool — guidelines checker spanning full width at bottom
    add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5),
             RGBColor(0x10, 0x20, 0x38))
    add_text(slide,
             "🛡  content_guidelines_checker — Deterministic brand-safety validator: "
             "detects prohibited content, over-promotion & platform violations. "
             "Runs alongside the LLM reviewer in the Check node.",
             Inches(0.5), Inches(6.82), Inches(12.3), Inches(0.44),
             font_size=11, color=ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Prompt Caching
# ═══════════════════════════════════════════════════════════════════════════════
def slide_caching():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "OPTIMISATION — PROMPT CACHING")
    slide_title(slide, "Prompt Caching — S9 Pattern  ⚡",
                "Azure OpenAI automatic prefix caching — reduces cost & latency on repeated calls")

    # Left: how it works
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(5.9), Inches(5.2),
             RGBColor(0x0A, 0x18, 0x2A))
    add_text(slide, "How it works", Inches(0.5), Inches(2.05), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=ACCENT)
    how = [
        "Every invoke_cached() call sends 2 messages:",
        "  SystemMessage: _CACHE_ANCHOR + agent-specific prompt",
        "  HumanMessage : dynamic per-request content",
        "_CACHE_ANCHOR is ~1 100 tokens — identical for ALL agents",
        "Azure caches automatically when same prefix ≥ 1 024 tokens",
        "Cache HIT from the 2nd request onward (same session)",
        "prompt_cache_key = 'langchain-prompt-caching' signals Azure",
        "API version 2024-12-01-preview required",
    ]
    add_bullet_box(slide, how, Inches(0.5), Inches(2.5), Inches(5.6), Inches(3.5),
                   font_size=12, bullet_color=WHITE)

    add_text(slide, "Token savings per generation", Inches(0.5), Inches(6.1), Inches(5.7), Inches(0.35),
             font_size=13, bold=True, color=ACCENT)
    add_text(slide,
             "8 LLM calls × ~1 100 cached tokens = 8 800 cached tokens/run\n"
             "Cached tokens billed at ~50% vs standard prompt tokens",
             Inches(0.5), Inches(6.45), Inches(5.6), Inches(0.6),
             font_size=12, color=LIGHT_GRAY)

    # Right: code snippet
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.2),
             RGBColor(0x06, 0x12, 0x06))
    add_text(slide, "invoke_cached() helper", Inches(6.9), Inches(2.05), Inches(5.9), Inches(0.4),
             font_size=14, bold=True, color=GREEN)
    code = """def invoke_cached(system_text, human_text, logger) -> LLMResult:
    messages = [
        SystemMessage(
            content=_CACHE_ANCHOR + system_text  # ← static prefix
        ),
        HumanMessage(content=human_text),  # ← dynamic
    ]
    response = llm.invoke(messages)

    cached_tokens = response.response_metadata\\
        ["token_usage"]["prompt_tokens_details"]\\
        ["cached_tokens"]

    return LLMResult(
        content=str(response.content),
        cache_status="HIT" if cached_tokens > 0 else "MISS",
        usage={...},  # prompt/completion/total/cached
        message=response,
    )"""
    add_text(slide, code, Inches(6.9), Inches(2.5), Inches(5.8), Inches(4.3),
             font_size=9.5, color=GREEN)

    # Response example
    add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.25),
             RGBColor(0x10, 0x20, 0x38))
    add_text(slide,
             "Response includes: token_usage { prompt_tokens, completion_tokens, "
             "total_tokens, cached_tokens, cache_hits, cache_misses } — returned in meta field of /generate response",
             Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.22),
             font_size=9, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — API Endpoints
# ═══════════════════════════════════════════════════════════════════════════════
def slide_api():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "API")
    slide_title(slide, "FastAPI Endpoints  🚀",
                "Two REST endpoints — upload source material then generate content")

    # Upload endpoint
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(12.5), Inches(2.3),
             RGBColor(0x0A, 0x20, 0x15))
    add_rect(slide, Inches(0.4), Inches(2.0), Inches(1.3), Inches(2.3),
             RGBColor(0x00, 0x88, 0x44))
    add_text(slide, "POST", Inches(0.45), Inches(2.7), Inches(1.2), Inches(0.5),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "/upload", Inches(1.8), Inches(2.05), Inches(4), Inches(0.5),
             font_size=18, bold=True, color=GREEN)
    add_text(slide, "multipart/form-data  •  field: file",
             Inches(1.8), Inches(2.55), Inches(6), Inches(0.35),
             font_size=12, color=LIGHT_GRAY)
    add_text(slide,
             "Accepts .txt .md .rst .csv .pdf  •  Saves with random hex name  •  "
             'Returns: { "file_path": "uploads/a3f...", "original_filename": "post.pdf" }',
             Inches(1.8), Inches(2.9), Inches(10.9), Inches(0.45),
             font_size=11, color=WHITE)
    add_text(slide, "Pass file_path as uploaded_file_path in /generate request",
             Inches(1.8), Inches(3.4), Inches(10.9), Inches(0.35),
             font_size=11, color=YELLOW, italic=True)

    # Generate endpoint
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.7),
             RGBColor(0x0A, 0x18, 0x2A))
    add_rect(slide, Inches(0.4), Inches(4.5), Inches(1.3), Inches(2.7),
             RGBColor(0x00, 0x55, 0xAA))
    add_text(slide, "POST", Inches(0.45), Inches(5.5), Inches(1.2), Inches(0.5),
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "/generate", Inches(1.8), Inches(4.55), Inches(4), Inches(0.5),
             font_size=18, bold=True, color=ACCENT)
    add_text(slide, "application/json  •  ContentRequest body",
             Inches(1.8), Inches(5.05), Inches(6), Inches(0.35),
             font_size=12, color=LIGHT_GRAY)

    req = ('{\n  "input_content": "...",   // optional\n'
           '  "source_url": "https://...", // optional\n'
           '  "uploaded_file_path": "...", // optional  ← at least one required\n'
           '  "max_iterations": 3          // default\n}')
    add_text(slide, "Request:", Inches(1.8), Inches(5.45), Inches(2), Inches(0.3),
             font_size=11, color=LIGHT_GRAY, bold=True)
    add_text(slide, req, Inches(1.8), Inches(5.75), Inches(5.2), Inches(1.3),
             font_size=9.5, color=WHITE)

    resp = ('{\n  "linkedin":  { "content": "...", "status": "ready" },\n'
            '  "twitter":   { "content": "...", "status": "ready" },\n'
            '  "instagram": { "content": "...", "status": "ready" },\n'
            '  "meta": { "iterations": 2, "overall_score": 8.9,\n'
            '            "approved": true,\n'
            '            "token_usage": { "prompt_tokens": 4800,\n'
            '              "cached_tokens": 3200, "cache_hits": 5 } }\n}')
    add_text(slide, "Response:", Inches(7.2), Inches(5.45), Inches(2), Inches(0.3),
             font_size=11, color=LIGHT_GRAY, bold=True)
    add_text(slide, resp, Inches(7.2), Inches(5.75), Inches(5.5), Inches(1.45),
             font_size=9.5, color=GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Key Innovations & Challenges
# ═══════════════════════════════════════════════════════════════════════════════
def slide_innovations():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)
    section_header(slide, "KEY INNOVATIONS")
    slide_title(slide, "Key Innovations & Challenges Solved  💡")

    innovations = [
        ("🔄  Targeted Reflection Loop",
         ACCENT,
         ["Standard reflect = rerun ALL agents every iteration",
          "Our approach: only retry platforms scoring < 7.0",
          "Prevents oscillation where passing platforms regress",
          "Saves 2–4 LLM calls per iteration in multi-platform setup"]),
        ("⚡  Prompt Caching Architecture",
         GREEN,
         ["Shared 1 100-token agency brief as common prefix",
          "All 8 calls/run share the same cacheable prefix",
          "Cache HIT from second request = ~50% cost reduction",
          "LLMResult dataclass surfaces cache_status per call"]),
        ("📊  Score Stagnation Guard",
         YELLOW,
         ["Problem: LLM reviewers have a scoring ceiling (~8.7)",
          "Without guard: 3 iterations all score 8.7, all waste tokens",
          "Fix: exit loop if score didn't improve vs previous iteration",
          "Combined with lowered threshold (8.8) prevents wasted loops"]),
        ("📎  Multi-Source Input",
         ORANGE,
         ["Any combination of: text brief / URL / uploaded file",
          "PDF extraction (pypdf), HTML stripping (stdlib urllib)",
          "All three source fields optional — at least one required",
          "Error-string guard prevents tool errors from poisoning plan"]),
    ]

    for i, (title, color, points) in enumerate(innovations):
        row = i // 2
        col = i % 2
        x = Inches(0.4 + col * 6.5)
        y = Inches(2.0 + row * 2.75)
        add_rect(slide, x, y, Inches(6.0), Inches(2.55), RGBColor(0x10, 0x20, 0x38))
        add_rect(slide, x, y, Inches(6.0), Inches(0.08), color)
        add_text(slide, title, x + Inches(0.1), y + Inches(0.12), Inches(5.8), Inches(0.45),
                 font_size=14, bold=True, color=color)
        add_bullet_box(slide, points, x + Inches(0.1), y + Inches(0.6), Inches(5.8), Inches(1.9),
                       font_size=12, bullet_color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Summary / Thank You
# ═══════════════════════════════════════════════════════════════════════════════
def slide_summary():
    slide = prs.slides.add_slide(blank_layout)
    fill_bg(slide)

    # Top gradient bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), ACCENT)

    add_text(slide, "Thank You  🙌", Inches(0.4), Inches(0.4), Inches(12.5), Inches(0.8),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Autonomous Social Media Content Studio",
             Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.45),
             font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # Summary grid
    summary = [
        ("Multi-Agent Pattern", "Plan-Act-Check-Reflect (LangGraph state machine)",   ACCENT),
        ("Parallelism",         "3 platform agents run concurrently via fan-out",      GREEN),
        ("Self-correction",     "Targeted reflect loop — only failing agents retry",   ACCENT2),
        ("Source flexibility",  "Text  /  URL  /  Uploaded file (PDF, TXT, MD…)",      YELLOW),
        ("Quality control",     "LLM scorer + deterministic brand-safety checker",     ORANGE),
        ("Cost optimisation",   "Prompt caching — shared 1 100-token static prefix",  ACCENT),
        ("Observability",       "token_usage in every response (HIT/MISS per call)",  LIGHT_GRAY),
        ("API",                 "FastAPI — /upload  +  /generate endpoints",          GREEN),
    ]

    for i, (key, val, col) in enumerate(summary):
        row = i // 2
        c = i % 2
        x = Inches(0.4 + c * 6.5)
        y = Inches(2.0 + row * 1.1)
        add_rect(slide, x, y, Inches(6.0), Inches(0.95), RGBColor(0x10, 0x20, 0x38))
        add_rect(slide, x, y, Inches(6.0), Inches(0.06), col)
        add_text(slide, key, x + Inches(0.1), y + Inches(0.08), Inches(2.2), Inches(0.4),
                 font_size=12, bold=True, color=col)
        add_text(slide, val, x + Inches(0.1), y + Inches(0.47), Inches(5.8), Inches(0.44),
                 font_size=11, color=WHITE)

    add_rect(slide, Inches(0), Inches(6.6), SLIDE_W, Inches(0.9), RGBColor(0x10, 0x20, 0x38))
    add_text(slide,
             "Contributors: Himanshu Jain  •  Prashanchal Sharma                    "
             "Built with: LangGraph  •  LangChain  •  Azure OpenAI  •  FastAPI  •  Python 3.13",
             Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.75),
             font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Build all slides ──────────────────────────────────────────────────────────
slide_cover()
slide_problem_solution()
slide_strategy()
slide_architecture()
slide_workflow()
slide_state()
slide_planner()
slide_act()
slide_check()
slide_reflect()
slide_tools()
slide_caching()
slide_api()
slide_innovations()
slide_summary()

out_path = "demo_presentation_v2.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
